#!/usr/bin/env python3
"""Generate McKinsey-style GW SDLC Accelerator slide pack (.pptx).

Each of the 7 AI accelerators gets a dedicated slide with:
  - Action title + so-what bar (McKinsey convention)
  - Purpose, worked example, productionalisation path
  - Demo screenshot or GIF (right column)
  - ROI metric card
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DEMO_DIR = ROOT / "docs" / "media"
GIF_DIR = ROOT / "public" / "media"

# McKinsey-inspired palette
NAVY = RGBColor(0x05, 0x1C, 0x2C)
NAVY_MID = RGBColor(0x0D, 0x2D, 0x45)
TEAL = RGBColor(0x00, 0xA9, 0xCE)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_100 = RGBColor(0xF4, 0xF6, 0xF8)
GRAY_400 = RGBColor(0x8A, 0x9B, 0xAB)
GRAY_600 = RGBColor(0x5A, 0x6D, 0x7E)
TEXT = RGBColor(0x1A, 0x2B, 0x3C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUT = ROOT / "docs" / "accelerators-slide-pack.pptx"

ACCELERATORS = [
    {
        "id": "story-forge",
        "name": "Story Forge",
        "phase": "Plan",
        "action": "Story Forge converts raw requirements into sprint-ready, Guidewire-mapped user stories in minutes",
        "so_what": "Requirements in, sprint-ready stories out — Gherkin ACs, GW touchpoints and Fibonacci points pre-populated; fewer AC gaps reach SIT.",
        "purpose": "Eliminates the blank-page problem for BAs. Produces INVEST stories with negative-path ACs, mapped entities/PCF/plugins, and Jira-ready JSON.",
        "example": "Input: 'Apply multi-vehicle discount on mid-term PA change.' → Output: ST-1 with Given/When/Then, touchpoints (PolicyLine, RatingPlugin), 5 pts, open question on eligibility.",
        "prod": [
            "Pilot on live backlog grooming; map JSON → Jira/ADO fields",
            "Load client story standards into Project knowledge (RAG)",
            "Add SSO + tenant scoping before production client data",
        ],
        "roi": "40–60%",
        "roi_label": "Less effort per story",
    },
    {
        "id": "code-review",
        "name": "Code Review Copilot",
        "phase": "Build",
        "action": "Code Review Copilot delivers principal-level Guidewire review on every commit — catching defects at desk-check cost",
        "so_what": "Defects caught at review cost ~10× less than at UAT — zero senior reviewer calendar time; bundled GW Cloud standards per profile.",
        "purpose": "Reviews Gosu, PCF, integration, GX and batch code. Severity findings (critical→info), upgrade-safety flags, code-health score, optional Sonar merge.",
        "example": "Input: Gosu with gw.transaction in PCF onClick → Output: critical (UI-layer transaction), major (query in loop), score 62, concrete fix naming GW construct.",
        "prod": [
            "Pre-PR gate in CI pipeline; persist findings as artifacts",
            "Index client coding standards into Project knowledge",
            "Golden eval set in CI on prompt/model changes",
        ],
        "roi": "30%",
        "roi_label": "Fewer review cycles",
    },
    {
        "id": "test-strategist",
        "name": "Test Strategist",
        "phase": "Test",
        "action": "Test Strategist enforces pyramid discipline — assigning GUnit, GT-API, GT-UI or Manual per case",
        "so_what": "Test design keeps pace with development — harness selection enforced, test data staged upfront, full AC traceability.",
        "purpose": "Derives executable cases from stories, code or defects. Maintains pyramid health; lists preconditions and data to stage before execution.",
        "example": "Input: FNOL user story + ACs → Output: TC-1 GT-UI journey, TC-2 GUnit reserve rules, TC-3 GT-API claim create; data: in-force PA policy, vehicle VIN.",
        "prod": [
            "Webhook from ALM to auto-feed sprint stories",
            "Export to Zephyr/qTest via CSV transform",
            "Pair GT-UI cases with Flow Automator scripts",
        ],
        "roi": "50%",
        "roi_label": "Faster test-case design",
    },
    {
        "id": "flow-automator",
        "name": "Flow Automator",
        "phase": "Test",
        "action": "Flow Automator scaffolds keyword-driven Katalon UI automation for common Guidewire journeys in minutes, not days",
        "so_what": "Regression UI packs in minutes — real /katalon @Keyword signatures injected into prompts, not guessed method names.",
        "purpose": "Generates Katalon Groovy for PC/CC/BC/Jutro flows (submission→bind, FNOL, billing, quote-and-buy). Matches bundled keyword libraries.",
        "example": "Input: PC — 'PA submission → quote → bind' → Output: Groovy calling createPersonAccount, startSubmission, quote, issuePolicy + prerequisites list.",
        "prod": [
            "Open /katalon in Katalon Studio; adapt locators per env",
            "Run katalonc headless in CI with profile secrets",
            "Regenerate after PCF customisation via Flow Automator",
        ],
        "roi": "60%+",
        "roi_label": "Faster UI automation draft",
    },
    {
        "id": "test-migrator",
        "name": "Test Migrator",
        "phase": "Test",
        "action": "Test Migrator converts legacy manual regression packs into runnable automation — with honest gap analysis per case",
        "so_what": "Manual pack → automation drafts in minutes — per-case verdict (automate / fix-then-automate / keep manual) plus test-data manifest.",
        "purpose": "Ingests Excel/Zephyr/qTest cases (single or bulk). Converts to Katalon, GT, Playwright, Selenium or Cucumber; surfaces blocking gaps.",
        "example": "Input: TC-PC-014 manual bind from Excel → Output: Katalon script, verdict automate-with-fixes, gap (missing underwriter role), data: person account (generate).",
        "prod": [
            "Batch-convert client UAT pack in pilot weeks 2–4",
            "Remediate gaps before promoting scripts to CI",
            "Wire Katalon output into /katalon regression suite",
        ],
        "roi": "70%+",
        "roi_label": "Faster manual conversion",
    },
    {
        "id": "release-navigator",
        "name": "Release Navigator",
        "phase": "Release",
        "action": "Release Navigator reveals what the next ski release does to your customisations — before it lands",
        "so_what": "Upgrade planning from evidence, not guesswork — CI/CD maturity score + per-area impact, remediation and regression focus.",
        "purpose": "16-practice readiness self-check plus AI impact analysis of customisation inventory vs target ski release; effort band and pre-upgrade checklist.",
        "example": "Input: Palisades, PC+CC, inventory (IG, entity extensions, PCF) → Output: medium risk, high IG impact, 8-step checklist, targeted regression list.",
        "prod": [
            "Maintain inventory in Project knowledge; re-run each ski release",
            "Attach official GW release notes to RAG corpus",
            "Feed regression focus into Test Strategist + Katalon",
        ],
        "roi": "20–30%",
        "roi_label": "Less upgrade regression effort",
    },
    {
        "id": "defect-triage",
        "name": "Defect Triage Agent",
        "phase": "Operate",
        "action": "Defect Triage Agent runs a four-agent pipeline with self-correction — L2/L3 triage in minutes, 24/7",
        "so_what": "The flagship agentic module — confidence-gated loop (max 2 passes) with live timeline; ideal for client showcases.",
        "purpose": "Intake → Investigator → Router (loops if confidence < 65%) → Fix Planner. Datadog log integration; bulk triage up to 5 errors.",
        "example": "Input: Payment posting failure + Datadog logs → Output: P2 Billing, 72% confidence hypothesis, workaround, permanent fix, regression cases.",
        "prod": [
            "Connect DD_API_KEY for live log evidence",
            "Index client runbooks in Project knowledge",
            "Push case file JSON to ServiceNow / Jira enrichment",
        ],
        "roi": "35–50%",
        "roi_label": "Faster mean-time-to-triage",
    },
]


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=TEXT,
                align=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_labeled_block(slide, left, top, width, label, body, label_size=9, body_size=10):
    add_textbox(slide, left, top, width, Inches(0.22), label, size=label_size, bold=True, color=TEAL)
    add_textbox(slide, left, top + Inches(0.2), width, Inches(0.55), body, size=body_size, color=TEXT)
    return top + Inches(0.78)


def add_bullets(slide, left, top, width, height, items, size=11, color=TEXT, spacing=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(spacing)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.bullet = True
    return box


def add_footer(slide, page_num):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(3), Inches(0.3),
                "GW SDLC Accelerator", size=8, bold=True, color=NAVY)
    add_textbox(slide, Inches(5.5), Inches(7.05), Inches(2.3), Inches(0.3),
                "Confidential", size=8, color=GRAY_400, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(12.3), Inches(7.05), Inches(0.5), Inches(0.3),
                str(page_num), size=8, color=GRAY_400, align=PP_ALIGN.RIGHT)
    fill_rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.01), GRAY_100)


def add_header_slide(slide, action_title, subtitle):
    fill_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.35), NAVY)
    add_textbox(slide, Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.75),
                action_title, size=20, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.3),
                subtitle, size=10, color=GRAY_400)


def add_so_what(slide, text, top=Inches(1.55)):
    fill_rect(slide, Inches(0.5), top, Inches(12.3), Inches(0.65), GRAY_100)
    fill_rect(slide, Inches(0.5), top, Inches(0.06), Inches(0.65), TEAL)
    add_textbox(slide, Inches(0.7), top + Inches(0.08), Inches(12), Inches(0.5),
                text, size=11, bold=True, color=NAVY)


def add_metric_card(slide, left, top, value, label, width=Inches(2.8)):
    fill_rect(slide, left, top, width, Inches(1.05), GRAY_100)
    fill_rect(slide, left, top, width, Inches(0.05), TEAL)
    add_textbox(slide, left, top + Inches(0.12), width, Inches(0.5),
                value, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.58), width - Inches(0.2), Inches(0.4),
                label, size=8, color=GRAY_600, align=PP_ALIGN.CENTER)


def add_table(slide, left, top, width, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.32 * n_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_100
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = TEXT
                if j == 0:
                    p.font.bold = True
    return table_shape


def demo_path(module_id: str) -> Path | None:
    gif = GIF_DIR / f"{module_id}.gif"
    png = DEMO_DIR / f"{module_id}-demo.png"
    if gif.exists():
        return gif
    if png.exists():
        return png
    return None


def add_demo_image(slide, module_id: str):
    path = demo_path(module_id)
    if not path:
        fill_rect(slide, Inches(7.0), Inches(2.35), Inches(5.8), Inches(4.35), GRAY_100)
        add_textbox(slide, Inches(7.2), Inches(4.2), Inches(5.4), Inches(0.5),
                    "Demo preview\n(run npm run slides to generate)", size=10, color=GRAY_400, align=PP_ALIGN.CENTER)
        return
    # McKinsey layout: demo frame with teal top rule
    fill_rect(slide, Inches(7.0), Inches(2.35), Inches(5.8), Inches(4.35), GRAY_100)
    fill_rect(slide, Inches(7.0), Inches(2.35), Inches(5.8), Inches(0.06), TEAL)
    add_textbox(slide, Inches(7.15), Inches(2.42), Inches(3), Inches(0.2),
                "LIVE DEMO", size=7, bold=True, color=TEAL)
    slide.shapes.add_picture(str(path), Inches(7.15), Inches(2.65), Inches(5.5), Inches(3.9))


def slide_accelerator(prs, acc: dict, page: int) -> int:
    slide = blank_slide(prs)
    add_header_slide(slide, acc["action"], f"{acc['phase']} · {acc['name']}")
    add_so_what(slide, acc["so_what"])

    y = Inches(2.35)
    y = add_labeled_block(slide, Inches(0.5), y, Inches(6.2), "PURPOSE", acc["purpose"])
    y = add_labeled_block(slide, Inches(0.5), y, Inches(6.2), "EXAMPLE", acc["example"], body_size=9)
    add_textbox(slide, Inches(0.5), y, Inches(2), Inches(0.22), "PRODUCTIONALISE", size=9, bold=True, color=TEAL)
    add_bullets(slide, Inches(0.5), y + Inches(0.22), Inches(6.2), Inches(1.5), acc["prod"], size=9, spacing=4)

    add_metric_card(slide, Inches(0.5), Inches(6.15), acc["roi"], acc["roi_label"], width=Inches(2.4))
    add_demo_image(slide, acc["id"])
    add_footer(slide, page)
    return page


def ensure_demos():
    demo_script = ROOT / "scripts" / "generate-accelerator-demos.py"
    if demo_script.exists():
        subprocess.run([sys.executable, str(demo_script)], check=True)
    for gif_script in ["generate-code-review-gif.py", "generate-defect-triage-gif.py"]:
        path = ROOT / "scripts" / gif_script
        if path.exists():
            subprocess.run([sys.executable, str(path)], check=False)


def slide_title(prs):
    slide = blank_slide(prs)
    fill_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    fill_rect(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), TEAL)
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(0.4),
                "NTT DATA GUIDEWIRE PRACTICE", size=10, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(1.2),
                "GuidewireAI\nSDLC Test Accelerators", size=40, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(9), Inches(0.8),
                "Seven AI modules for Guidewire InsuranceSuite — purpose, demo, examples & productionalisation",
                size=16, color=GRAY_400)
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(10), Inches(0.6),
                "McKinsey-style client pack · July 2026 · Confidential",
                size=10, color=GRAY_400)


def slide_section(prs, num, title, desc):
    slide = blank_slide(prs)
    fill_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY_MID)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(2), Inches(1.2), num, size=72, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(0.8), title, size=32, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.4), Inches(8), Inches(0.8), desc, size=14, color=GRAY_400)


def slide_next_steps(prs):
    slide = blank_slide(prs)
    fill_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    fill_rect(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), TEAL)
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(4), Inches(0.4), "NEXT STEPS", size=10, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(0.8),
                "Deploy. Demo. Deliver.", size=36, color=WHITE)
    add_bullets(slide, Inches(0.8), Inches(3.3), Inches(9), Inches(2.2), [
        "Deploy to Vercel — Anthropic + Neon + optional Voyage for RAG embeddings",
        "Schedule client showcase: Defect Triage + Flow Automator + Katalon live",
        "Pilot weeks 2–4 on live client data; measure ROI per module",
        "Harden: SSO, rate limits, CI evals — see productionalisation roadmap",
    ], size=13, color=GRAY_400, spacing=10)
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(10), Inches(0.6),
                "NTT DATA Guidewire Practice · GuidewireAI · gw-sdlc-accelerator",
                size=11, color=TEAL)


def build():
    ensure_demos()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    page = 0

    slide_title(prs)

    # Executive summary
    slide = blank_slide(prs)
    add_header_slide(slide,
        "Seven AI accelerators compress manual work across the Guidewire testing lifecycle — from stories to triage",
        "Executive Summary")
    add_so_what(slide,
        "GuidewireAI packages Claude-powered modules, bundled reference corpora (/katalon, GW standards), and per-project RAG into one web application.")
    add_textbox(slide, Inches(0.5), Inches(2.35), Inches(2), Inches(0.25), "THE CHALLENGE", size=9, bold=True, color=TEAL)
    add_bullets(slide, Inches(0.5), Inches(2.6), Inches(5.8), Inches(2.2), [
        "Test design, UI automation, manual-to-auto conversion and triage consume 40–60% of sprint capacity",
        "Inconsistent standards across delivery centres; UI regression rots with every customisation",
        "AMS teams face 24/7 defect volume with limited L2/L3 capacity",
    ], size=10)
    add_textbox(slide, Inches(6.8), Inches(2.35), Inches(2), Inches(0.25), "OUR RESPONSE", size=9, bold=True, color=TEAL)
    add_bullets(slide, Inches(6.8), Inches(2.6), Inches(5.8), Inches(2.2), [
        "7 AI modules: Plan → Build → Test → Release → Operate",
        "Bundled reference material + per-project knowledge (RAG)",
        "Katalon companion: 12 ready-to-run flows (PC · CC · BC · Jutro)",
        "Indicative ROI: 30–70% effort reduction on targeted activities",
    ], size=10)
    add_metric_card(slide, Inches(0.5), Inches(5.1), "7", "AI accelerators")
    add_metric_card(slide, Inches(3.6), Inches(5.1), "12", "Katalon flows")
    add_metric_card(slide, Inches(6.7), Inches(5.1), "5", "SDLC phases")
    add_metric_card(slide, Inches(9.8), Inches(5.1), "RAG", "Per-project knowledge")
    add_footer(slide, page := 2)

    # Lifecycle
    slide = blank_slide(prs)
    add_header_slide(slide,
        "Each SDLC phase has a dedicated accelerator — a connected value chain from backlog to production",
        "Lifecycle Architecture")
    phases = [
        ("PLAN", "Story Forge"),
        ("BUILD", "Code Review"),
        ("TEST", "Test Strategist · Flow Automator · Test Migrator"),
        ("RELEASE", "Release Navigator"),
        ("OPERATE", "Defect Triage Agent"),
    ]
    x = Inches(0.4)
    w = Inches(2.45)
    for label, box in phases:
        add_textbox(slide, x, Inches(1.55), w, Inches(0.25), label, size=7, bold=True, color=GRAY_400, align=PP_ALIGN.CENTER)
        fill_rect(slide, x, Inches(1.85), w, Inches(0.55), NAVY)
        add_textbox(slide, x + Inches(0.05), Inches(1.9), w - Inches(0.1), Inches(0.45),
                    box, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += w + Inches(0.12)
    add_table(slide, Inches(0.5), Inches(2.7), Inches(12.3),
              ["Accelerator", "Phase", "Primary testing value"],
              [[a["name"], a["phase"], a["roi_label"]] for a in ACCELERATORS],
              [Inches(2.2), Inches(1.0), Inches(9.1)])
    add_footer(slide, page := 3)

    slide_section(prs, "01", "The seven AI accelerators",
                  "Each slide: purpose · worked example · demo preview · productionalisation path")

    page = 3
    for acc in ACCELERATORS:
        page = slide_accelerator(prs, acc, page + 1)

    slide_section(prs, "02", "Productionalise the suite",
                  "From internal demo to client-ready, governed enterprise service")

    # Production roadmap
    slide = blank_slide(prs)
    add_header_slide(slide,
        "A four-phase roadmap moves GuidewireAI from demo to enterprise scale without re-architecting",
        "Productionalisation Roadmap")
    add_so_what(slide, "Phase 0 works today on Vercel + Anthropic + Neon. Phases 1–3 close auth, tenancy, CI evals and Katalon automation gates.")
    add_table(slide, Inches(0.5), Inches(2.35), Inches(12.3),
              ["Phase", "Goal", "Key work"],
              [
                  ["0 — Now", "Internal / client demo", "Vercel deploy; bundled reference + RAG; /katalon open in Studio"],
                  ["1 — Pilot", "Named users, safe pilot", "SSO, rate limits, input validation, DB migrations, CI build gate"],
                  ["2 — Client-ready", "Multi-tenant, governed", "Row-level tenancy, audit log, prompt evals in CI, Katalon in CI (KRE)"],
                  ["3 — Enterprise", "Supportable product", "Token budgets, schema-validated outputs, TestOps, API-driven test data"],
              ],
              [Inches(1.0), Inches(2.2), Inches(9.1)])
    add_footer(slide, page := page + 1)

    # Katalon companion
    slide = blank_slide(prs)
    add_header_slide(slide,
        "Katalon Flow Automation ships 12 ready-to-run flows — the execution layer behind Flow Automator",
        "Companion Asset · /katalon")
    add_so_what(slide, "AI generates scripts; Katalon runs them. Keyword libraries centralise locators — one fix per customised screen.")
    add_table(slide, Inches(0.5), Inches(2.35), Inches(6.2),
              ["Product", "Flows"],
              [
                  ["PolicyCenter", "Submission→bind · Change · Cancel · Renewal"],
                  ["ClaimCenter", "FNOL · Reserve+payment · Close"],
                  ["BillingCenter", "Payment · Invoices · Disbursement"],
                  ["Jutro", "Quote-and-buy · Self-service FNOL"],
              ],
              [Inches(1.4), Inches(4.8)])
    add_bullets(slide, Inches(7), Inches(2.35), Inches(5.8), Inches(2.5), [
        "Productionalise: katalonc headless in CI",
        "Secret-injected profiles (default.glbl, qa.glbl)",
        "Flow Automator injects real Groovy signatures from /katalon",
    ], size=10)
    add_footer(slide, page := page + 1)

    # Architecture
    slide = blank_slide(prs)
    add_header_slide(slide,
        "Serverless architecture keeps secrets server-side; bundled corpora and RAG enrich every accelerator run",
        "Technical Architecture")
    add_bullets(slide, Inches(0.5), Inches(2.35), Inches(6), Inches(3.5), [
        "React SPA → Vercel /api/chat (Edge) + /api/knowledge + /api/projects",
        "Claude Sonnet 4.6 — prompt caching on reference material",
        "Global corpora: /katalon, /reference (GW standards, ski releases)",
        "Per-project RAG: chunk → embed → retrieve top-8 chunks",
        "Neon Postgres: projects, artifacts, knowledge_docs/chunks",
        "Optional pgvector + Voyage/OpenAI embeddings at scale",
    ], size=10)
    add_bullets(slide, Inches(7), Inches(2.35), Inches(5.8), Inches(3), [
        "ANTHROPIC_API_KEY — server-side only",
        "DATABASE_URL — Neon pooled connection",
        "VOYAGE_API_KEY — production RAG embeddings (optional)",
        "DD_API_KEY — Defect Triage Datadog integration",
    ], size=10)
    add_footer(slide, page := page + 1)

    # ROI
    slide = blank_slide(prs)
    add_header_slide(slide,
        "Indicative ROI: 30–70% effort reduction concentrated at highest-volume manual activities",
        "Value at a Glance")
    add_table(slide, Inches(0.4), Inches(1.55), Inches(12.5),
              ["Accelerator", "Phase", "ROI", "Productionalise first step"],
              [[a["name"], a["phase"], f"{a['roi']} {a['roi_label']}", a["prod"][0]] for a in ACCELERATORS],
              [Inches(2.0), Inches(0.8), Inches(2.5), Inches(7.2)])
    add_footer(slide, page := page + 1)

    # Adoption
    slide = blank_slide(prs)
    add_header_slide(slide,
        "Phased adoption moves teams from showcase to embedded workflow in four steps",
        "Recommended Adoption Path")
    steps = [
        ("Week 1 — Discover & demo", "Lead with Defect Triage + Flow Automator + Katalon hands-on."),
        ("Weeks 2–4 — Pilot on live data", "Story Forge grooming, Code Review on real Gosu, Test Migrator on 10–20 cases."),
        ("Month 2+ — Embed", "Jira/ADO export, Katalon in CI, Code Review as pre-PR gate, Project knowledge populated."),
        ("Ongoing — Measure", "Track ROI per module; extend /katalon libraries and RAG corpus per client."),
    ]
    y = Inches(1.7)
    for i, (title, desc) in enumerate(steps, 1):
        fill_rect(slide, Inches(0.5), y, Inches(0.35), Inches(0.35), TEAL)
        add_textbox(slide, Inches(0.5), y + Inches(0.02), Inches(0.35), Inches(0.3),
                    str(i), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.0), y, Inches(11.5), Inches(0.3), title, size=12, bold=True, color=NAVY)
        add_textbox(slide, Inches(1.0), y + Inches(0.32), Inches(11.5), Inches(0.55), desc, size=10, color=GRAY_600)
        y += Inches(1.15)
    add_footer(slide, page := page + 1)

    slide_next_steps(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
